#!/usr/bin/env python3
"""15-slide simple presentation — plain language for non-expert audience."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = "/run/media/touhid/SSD/batt_sense/battery_sens/battery_sens"
FIG  = os.path.join(ROOT, "data")
OUT  = "/run/media/touhid/SSD/batt_sense/conference/conference -A/presentation"

os.makedirs(OUT, exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY  = RGBColor(0x55, 0x55, 0x55)
BLUE  = RGBColor(0x34, 0x8A, 0xBD)

BLANK = prs.slide_layouts[6]
CONTENT = prs.slide_layouts[1]

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def textbox(slide, left, top, width, height, text, size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Calibri"; p.alignment = align
    return tf

def add(tf, text, size=20, bold=False, color=DARK, space=Pt(8), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Calibri"
    p.space_before = space; p.alignment = align
    return p

def title_slide(title, subtitle=""):
    sl = prs.slides.add_slide(BLANK); add_bg(sl)
    textbox(sl, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5), title, size=38, bold=True)
    if subtitle:
        textbox(sl, Inches(1), Inches(3.8), Inches(11.3), Inches(1.2), subtitle, size=18, color=GRAY)
    return sl

def bul_slide(title, lines, subtitle=""):
    sl = prs.slides.add_slide(BLANK); add_bg(sl)
    textbox(sl, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8), title, size=30, bold=True)
    # Separator
    sl.shapes.add_shape(1, Inches(0.7), Inches(1.3), Inches(11.9), Pt(2)).fill.solid()
    if subtitle:
        textbox(sl, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5), subtitle, size=16, color=GRAY)
    tf = textbox(sl, Inches(0.7), Inches(2.1), Inches(11.9), Inches(5.0), "", size=20)
    tf.paragraphs[0].text = ""
    for i, line in enumerate(lines):
        add(tf, line, size=20, bold=(i == 0), space=Pt(6))
    return sl

def fig_slide(title, fig_name, caption=""):
    sl = prs.slides.add_slide(BLANK); add_bg(sl)
    textbox(sl, Inches(0.7), Inches(0.3), Inches(11.9), Inches(0.7), title, size=30, bold=True)
    path = os.path.join(FIG, fig_name)
    if os.path.exists(path):
        sl.shapes.add_picture(path, Inches(1.0), Inches(1.2), width=Inches(11.3))
    if caption:
        textbox(sl, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.8), caption, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    return sl

# ── SLIDE 1: TITLE ──
title_slide(
    "Multi-Horizon Hazard Models for\nBattery Failure Prediction",
    "Can we predict battery failure across different chemistries?\n\n[Author Name] · [Affiliation]"
)

# ── SLIDE 2: THE PROBLEM ──
bul_slide("The Problem", [
    "Batteries fail over time — we want to predict failure before it happens",
    "This helps EVs, phones, and grid storage operate safely",
    "Previous work tested only ONE battery type (LCO laptop batteries)",
    "Question 1: Do different AI models give different results?",
    "Question 2: Can we train on one chemistry and predict on another?"
])

# ── SLIDE 3: WHAT DOES FAILURE MEAN ──
bul_slide("What Does \"Failure\" Mean?", [
    "We check two vital signs, like a doctor checking health:",
    "  1. State-of-Health (SOH) — battery capacity dropped below 80%",
    "  2. Voltage Sag — power drops suddenly before capacity fades",
    "Either one triggered = \"failure\" label",
    "Once triggered, stays triggered — batteries don't heal"
])

# ── SLIDE 4: THE BATTERIES WE TESTED ──
bul_slide("The Batteries We Tested", [
    "NASA 18650: 37 LCO cells (laptop chemistry), ~300 cycles each",
    "CALCE LCO/CX2: 7 LCO cells, 775–1952 cycles each",
    "Oxford LFP: 5 LFP cells (EV chemistry), ~300 cycles each",
    "LCO → test. LFP → test. Train on LCO, predict LFP? That's the key question."
])

# ── SLIDE 5: HOW WE TESTED ──
bul_slide("How We Tested", [
    "Three AI models: XGBoost, LightGBM, Random Forest",
    "Four prediction windows: 10, 20, 30, 50 cycles ahead",
    "Two calibration methods: isotonic (freehand) vs Platt (smooth curve)",
    "⚠ Critical: test with SOH feature AND without it"
])

# ── SLIDE 6: FINDING 1 ──
fig_slide("Finding 1: Works on Known Batteries",
          "Fig01_Within_Dataset_AUC.png",
          "All three models score AUC ≥ 0.85 on NASA and CALCE — the framework is reliable on familiar chemistries")

# ── SLIDE 7: FINDING 2 ──
fig_slide("Finding 2: Longer Windows = Easier",
          "Fig05_MultiHorizon_AUC.png",
          "Predicting 50 cycles ahead is easier than 10 — more time gives more data to work with")

# ── SLIDE 8: FINDING 3 — CALIBRATION ──
fig_slide("Finding 3: Getting Probabilities Right",
          "Fig02_Calibration_Comparison.png",
          "Platt calibration (solid lines) always beats isotonic (dashed). Especially important for CALCE.")

# ── SLIDE 9: WHY PLATT WINS ──
bul_slide("Why Platt Calibration is Better", [
    "Think of isotonic as drawing a freehand squiggle to fit the data:",
    "  • Follows every bump — works fine on small, clean datasets",
    "  • But overfits on messy, long-tailed data (like CALCE's 8733 cycles)",
    "Platt uses a smooth S-curve instead:",
    "  • Like taking a step back and seeing the big picture",
    "  • More robust — it doesn't chase noise",
    "CALCE Brier: isotonic 0.098 → Platt 0.069 (30% improvement)"
])

# ── SLIDE 10: CROSS-CHEM WITH SOH ──
fig_slide("Finding 4a: Looks Like It Works...",
          "Fig03_CrossChem_With_SOH.png",
          "With SOH, AUC hits 0.87–1.00 — looks like the model transfers perfectly from LCO to LFP")

# ── SLIDE 11: CROSS-CHEM WITHOUT SOH ──
fig_slide("Finding 4b: ...But It's a Trick",
          "Fig04_CrossChem_No_SOH.png",
          "Without SOH, AUC collapses to 0.51–0.54 — the model was 'cheating' using SOH as a shortcut")

# ── SLIDE 12: THE MECHANISM ──
bul_slide("The \"SOH Lookup Table\"", [
    "When the model has SOH, it learns this rule:",
    "  \"SOH = 85% → about 50 more cycles before failure\"",
    "It memorizes this from LCO training data",
    "LFP batteries happen to have similar SOH values...",
    "...so the model applies the same wrong rule to LFP",
    "Result: predictions look good (high AUC)", 
    "But: the model learned nothing about LFP itself — just reused a lookup table"
])

# ── SLIDE 13: KEY TAKEAWAY ──
sl = prs.slides.add_slide(BLANK); add_bg(sl)
textbox(sl, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8),
        "What This Means", size=30, bold=True)
sl.shapes.add_shape(1, Inches(0.7), Inches(1.3), Inches(11.9), Pt(2)).fill.solid()
textbox(sl, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.0),
        "The hazard framework works reliably\nwithin the same battery chemistry\n(AUC ≥ 0.85).",
        size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
textbox(sl, Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.5),
        "But cross-chemistry transfer is an unsolved problem.\nYou cannot train on LCO and deploy on LFP without\ngood features that actually transfer.",
        size=20, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 14: CAVEATS ──
bul_slide("Caveats", [
    "Oxford LFP: only 5 cells — small sample", 
    "Tested one direction only (LCO → LFP)",
    "Tree-based models only — deep learning might do better",
    "Voltage sag useless for LFP (flat voltage plateau)",
    "Published Brier scores could not be reproduced (code mismatch)"
])

# ── SLIDE 15: NEXT STEPS ──
bul_slide("Next Steps", [
    "Test on larger datasets with more battery types",
    "Design features that are \"chemistry-agnostic\"",
    "Try deep learning (LSTM, transformers)",
    "Test both directions: LCO→LFP AND LFP→LCO"
])

out_path = os.path.join(OUT, "presentation_simple.pptx")
prs.save(out_path)
print(f"Saved: {out_path} ({len(prs.slides)} slides)")
